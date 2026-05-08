import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))
from pptx import Presentation
from pptx.util import Pt


def build_slide(prs, layout_idx: int, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0


def build_table_slide(prs, title: str, rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    tbl = slide.shapes.add_table(
        len(rows), len(rows[0]), 914400, 1600000, 7315200, 2400000
    ).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val


# 31 content slides + 1 table slide = 32 (Phase 2 guideline deck)
SLIDES = [
    (
        0,
        "TraceDoc — Phase 2 Agentic Desk",
        [
            "Team: [Your names]",
            "CS-4063 · Professional Practices in IT",
            "Deliverable: Flask + Gemini Vision + HITL web desk",
        ],
    ),
    (
        1,
        "Phase 1 Recap (baseline)",
        [
            "Goal: English document images → formatted .docx",
            "Phase 1 stack: local preprocessing + OCR-style extraction",
            "Limitation: static pipeline, no autonomy, no explainability layer",
        ],
    ),
    (
        1,
        "Computing as a Formal Profession",
        [
            "Software work affects trust in legal, academic, and office records",
            "We owe competence, honesty, and proportionate safeguards to users",
            "TraceDoc treats every upload as potential sensitive data",
        ],
    ),
    (
        1,
        "Ethics vs Morals (applied)",
        [
            "Morals: ship fast, skip review",
            "Ethics: document risks, log reasoning, delete temps promptly",
            "UI + server enforce acknowledgement when confidence < 80%",
        ],
    ),
    (
        1,
        "Professional Ethics in Engineering",
        [
            "Security: path-safe downloads, no traversal on /api/download",
            "Quality: heuristic decide() scores JSON structure, not vibes",
            "User safety: abort path when extraction is structurally invalid",
        ],
    ),
    (
        1,
        "Ethical Decision (speed vs rigor)",
        [
            "Issue: Gemini is fast but can hallucinate layout JSON",
            "Choice: add structural heuristics + human gate before export",
            "Trade-off: slightly slower UX for higher accountability",
        ],
    ),
    (
        1,
        "Why Ethical Choices Matter Here",
        [
            "Wrong .docx can misrepresent contracts or coursework",
            "Society loses trust in “AI document” tools without transparency",
            "Immutable memory log supports post-hoc audit for demos",
        ],
    ),
    (
        1,
        "Ethical Theory + Human-Centered Design",
        [
            "Deontology: duty to delete temp images after act/abort",
            "Utilitarianism: minimize API calls via bounded retries/heuristics",
            "HCD: readable XAI stream + plain-language risk notes",
        ],
    ),
    (
        1,
        "ACM / IEEE (sample mapping)",
        [
            "ACM 1.2 Avoid harm: block reckless export when abort",
            "IEEE transparency: cite compliance_db keys in the XAI trace",
            "Tension: cloud inference — disclose in README + UI notice",
        ],
    ),
    (
        1,
        "4-Step Ethical Process (documented)",
        [
            "1. Identify risk: PII in uploads + model misuse",
            "2. Stakeholders: student user, document subjects, developers",
            "3. Alternatives: full cloud vs minimized retention (we chose latter)",
            "4. Decision: temp storage + immediate deletion post-session",
        ],
    ),
    (
        1,
        "Industry Practices vs Student Project",
        [
            "Industry: CI, threat modeling, SOC2-style logging",
            "Student repo: manual test + sample PDFs in TestSamples/",
            "Gap: formal pen-test; mitigate with input validation + least privilege keys",
        ],
    ),
    (
        1,
        "Trends: Agentic Document Systems",
        [
            "Shift from batch OCR to observe→decide→act loops",
            "LLMs as perception, rules as guardrails",
            "Regulators expect explainability alongside automation",
        ],
    ),
    (
        1,
        "Career Skills Demonstrated",
        [
            "Full-stack wiring: Flask REST + static SPA in /web",
            "Responsible AI: HITL gate + JSON risk scoring",
            "Communication: deck + screen recording walkthrough",
        ],
    ),
    (
        1,
        "Virtual Work + Sustainability",
        [
            "Remote-friendly: browser UI, no installer for reviewers",
            "Green angle: delete temps + avoid redundant Gemini calls on hard fail",
            "Honest claim: quantify vs Phase 1 in your narration",
        ],
    ),
    (
        1,
        "Legal Responsibilities (developer view)",
        [
            "PECA 2016 mindset: limit unauthorized processing of sensitive scans",
            "GDPR-style minimization: no long-term image retention in app",
            "User education: show compliance sidebar + data notice in UI",
        ],
    ),
    (
        1,
        "IPR + Licensing",
        [
            "Student-authored UI + agent glue code (set your license)",
            "python-docx / Flask / Google SDKs: follow upstream licenses",
            "Dataset: only use images you have rights to demonstrate",
        ],
    ),
    (
        1,
        "Computer Crime / Abuse Scenarios",
        [
            "Risk: adversarial text in image → prompt injection JSON",
            "Mitigation: injection substring scan collapses confidence",
            "Risk: malicious filename — server stores UUID.png only",
        ],
    ),
    (
        1,
        "Contracts + Terms (if you ship publicly)",
        [
            "Add ToS: no warranty on extracted text accuracy",
            "Age/consent if storing accounts (not in this MVP)",
            "Document override path for wrongful processing requests",
        ],
    ),
    (
        1,
        "Technical Limits of Phase 1",
        [
            "Single-shot conversion with limited feedback",
            "No structured explainability for graders or users",
            "No explicit autonomy level or memory story",
        ],
    ),
    (
        1,
        "Agentic System Concept",
        [
            "Perceive: preprocess + Gemini vision JSON",
            "Decide: heuristic confidence + policy checks",
            "Act: python-docx export after human authorization",
        ],
    ),
    (
        1,
        "Agentic Vision",
        [
            "Tool → Agent: system proposes export, human confirms",
            "Reactive → Proactive: streaming XAI trace before gate",
            "Learning: append-only JSON logs for audit demos",
        ],
    ),
    (
        1,
        "Architecture (repo layout)",
        [
            "implementation/web — SPA assets (HTML/CSS/JS)",
            "implementation/src — Flask app.py + agent/* + pipeline/*",
            "implementation/slides — programmatic PowerPoint builder",
        ],
    ),
    (
        1,
        "Agent Type",
        [
            "Goal-based agent: deliver accountable .docx with trace",
            "Bounded autonomy: server rejects act when decision is abort",
            "Justification: legal-grade docs need deterministic guardrails",
        ],
    ),
    (
        1,
        "Operational Workflow",
        [
            "POST /api/agent/perceive_decide with multipart image",
            "UI renders blocks + confidence gate",
            "POST /api/agent/act with optional acknowledge_risk under 80%",
        ],
    ),
    (
        1,
        "Intelligence Layer",
        [
            "LLM: Gemini vision for transcription + style hints",
            "Rules: decide() mean score + injection heuristics",
            "RAG-style law snippets: compliance_db.json via /api/compliance",
        ],
    ),
    (
        1,
        "Memory Design",
        [
            "Short-term: Flask SESSION_STATE dict per upload",
            "Long-term: confidence_log.json + xai_logs.json",
            "Nothing replaces enterprise KMS — scope honesty in talk track",
        ],
    ),
    (
        1,
        "Autonomy Level",
        [
            "Semi-autonomous: model proposes, human authorizes export",
            "Full autonomy deferred: would need stronger eval harness",
            "Demonstrate override via Abort session button",
        ],
    ),
    (
        1,
        "Human-in-the-Loop Controls",
        [
            "Gate screen shows computed % + risk copy",
            "Checkbox + JSON acknowledge_risk for <80% exports",
            "Download only serves sanitized .docx filenames",
        ],
    ),
    (
        1,
        "Ethical Agent Design",
        [
            "Privacy: temp file removed on success, abort, and error paths",
            "Bias: disclose English-print assumption in README",
            "Transparency: audit JSON column in UI after export",
        ],
    ),
    (
        1,
        "Risk Assessment",
        [
            "Model drift / API outage → empty JSON → abort UX",
            "Over-trust in green/yellow bars → require narration in video",
            "Misuse: do not run on confidential third-party data without consent",
        ],
    ),
    (
        1,
        "Safety Mechanisms",
        [
            "Structured logging to disk for demo replay",
            "Server-side enforcement mirrors UI warnings",
            "Easy kill-switch: revoke GEMINI_API_KEY in .env",
        ],
    ),
]

TABLE_SLIDE_DATA = [
    ["Feature", "Phase 1", "TraceDoc Phase 2"],
    ["Control", "User runs one script", "Agent proposes + human authorizes"],
    ["Intelligence", "Static OCR path", "LLM + heuristics + compliance RAG"],
    ["Behavior", "Reactive file-in/file-out", "Proactive XAI feed + policy gate"],
]

if __name__ == "__main__":
    prs = Presentation()
    for layout_idx, title, bullets in SLIDES[:20]:
        build_slide(prs, layout_idx, title, bullets)

    build_table_slide(prs, "Comparative analysis (guideline §III)", TABLE_SLIDE_DATA)

    for layout_idx, title, bullets in SLIDES[20:]:
        build_slide(prs, layout_idx, title, bullets)

    if len(prs.slides) != 32:
        raise ValueError(f"Slide count {len(prs.slides)} ≠ 32")

    out = Path(__file__).parent / "TraceDoc_Phase2_Deck.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"SLIDES GENERATED -> {out}")
